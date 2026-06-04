"""
うたかた SNS運用ご提案・第2回 スライド生成スクリプト
4:3 PowerPoint (.pptx) を生成
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ===== セットアップ =====
prs = Presentation()
prs.slide_width = Inches(10)     # 4:3
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ===== カラー =====
C_BG          = RGBColor(0xFA, 0xF7, 0xF2)
C_CARD        = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT        = RGBColor(0x2C, 0x2C, 0x2C)
C_PRIMARY     = RGBColor(0x5D, 0x7B, 0x7C)
C_PRIMARY_DK  = RGBColor(0x45, 0x61, 0x5F)
C_ACCENT      = RGBColor(0xB8, 0x91, 0x78)
C_ACCENT_DK   = RGBColor(0xA0, 0x7A, 0x62)
C_MUTED       = RGBColor(0x7A, 0x7A, 0x7A)
C_SOFT        = RGBColor(0xF5, 0xF1, 0xEA)
C_DEEP        = RGBColor(0x2D, 0x38, 0x38)
C_LINE        = RGBColor(0xE8, 0xE2, 0xD8)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_SUBTLE      = RGBColor(0xCF, 0xCB, 0xC0)

# ===== フォント =====
F_MINCHO  = "Yu Mincho"
F_GOTHIC  = "Yu Gothic"
F_MONO    = "Consolas"

# ===== ヘルパー =====
def add_blank_slide(bg=C_BG):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    return slide

def add_text(slide, x, y, w, h, text, *,
             font=F_MINCHO, size=14, color=C_TEXT,
             bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return tb

def add_rect(slide, x, y, w, h, *, fill=C_CARD, line=C_LINE, line_w=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape

def add_rounded(slide, x, y, w, h, *, fill=C_CARD, line=C_LINE, line_w=0.5, corner=0.04):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = corner
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape

def add_hline(slide, x, y, w, color=C_LINE, weight=0.5):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line

def add_header(slide, brand_text, page_num):
    """共通ヘッダー（ブランド名 + ページ番号 + 下線）"""
    add_text(slide, Inches(0.7), Inches(0.4), Inches(6), Inches(0.3),
             brand_text, font=F_MONO, size=9, color=C_MUTED, bold=True)
    add_text(slide, Inches(8.3), Inches(0.4), Inches(1.0), Inches(0.3),
             page_num, font=F_MONO, size=9, color=C_MUTED,
             align=PP_ALIGN.RIGHT)
    add_hline(slide, Inches(0.7), Inches(0.85), Inches(8.6))

def add_footer(slide, text, accent=False):
    """共通フッター（下線 + テキスト）"""
    add_hline(slide, Inches(0.7), Inches(6.85), Inches(8.6), color=C_LINE)
    color = C_ACCENT if accent else C_MUTED
    add_text(slide, Inches(0.7), Inches(6.95), Inches(8.6), Inches(0.3),
             text, font=F_GOTHIC, size=9, color=color,
             italic=not accent, bold=accent)

def add_title_block(slide, title, subtitle):
    """共通タイトル: 大見出し + サブタイトル"""
    add_text(slide, Inches(0.7), Inches(1.05), Inches(8.6), Inches(0.6),
             title, font=F_MINCHO, size=22, color=C_PRIMARY_DK, bold=False,
             line_spacing=1.3)
    add_text(slide, Inches(0.7), Inches(1.7), Inches(8.6), Inches(0.4),
             subtitle, font=F_GOTHIC, size=10, color=C_MUTED)

# ============================================================
# スライド 01: 表紙
# ============================================================
def slide_cover():
    slide = add_blank_slide(C_BG)
    # ブランドマーク
    add_text(slide, 0, Inches(2.2), SW, Inches(0.4),
             "U T A K A T A", font=F_MONO, size=11, color=C_MUTED,
             align=PP_ALIGN.CENTER, bold=True)
    # メインタイトル
    add_text(slide, 0, Inches(2.9), SW, Inches(0.9),
             "SNS運用ご提案", font=F_MINCHO, size=40, color=C_PRIMARY_DK,
             align=PP_ALIGN.CENTER, line_spacing=1.2)
    add_text(slide, 0, Inches(3.8), SW, Inches(0.9),
             "第2回", font=F_MINCHO, size=40, color=C_PRIMARY_DK,
             align=PP_ALIGN.CENTER, line_spacing=1.2)
    # 装飾ライン
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(4.7), Inches(5.0), Inches(0.6), Emu(9525))
    deco.fill.solid(); deco.fill.fore_color.rgb = C_ACCENT
    deco.line.fill.background()
    # サブ
    add_text(slide, 0, Inches(5.2), SW, Inches(0.5),
             "前回ヒアリングの振り返りと提案",
             font=F_MINCHO, size=13, color=C_TEXT, align=PP_ALIGN.CENTER)
    # メタ
    add_text(slide, 0, Inches(6.5), SW, Inches(0.3),
             "山田     /     2026",
             font=F_MONO, size=9, color=C_MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# スライド 02: アジェンダ
# ============================================================
def slide_agenda():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA", "02 / 12")
    add_title_block(slide, "アジェンダ",
                    "前回ヒアリングの振り返りと、今回の提案2点")

    items = [
        ("01", "前回ヒアリングの振り返り",  "前回出た課題の整理"),
        ("02", "提案① プロンプト一覧",     "投稿の言葉づくり用ツール"),
        ("03", "提案② 動画エディタ",       "テロップ位置調整不要の編集ツール"),
        ("04", "次回までのアクション",      "3週間の導入計画"),
    ]
    y = Inches(2.5)
    for num, title, desc in items:
        add_text(slide, Inches(0.9), y, Inches(0.7), Inches(0.4),
                 num, font=F_MONO, size=11, color=C_ACCENT, bold=True)
        add_text(slide, Inches(1.6), y - Inches(0.04), Inches(7.5), Inches(0.4),
                 title, font=F_MINCHO, size=15, color=C_TEXT, bold=True)
        add_text(slide, Inches(1.6), y + Inches(0.35), Inches(7.5), Inches(0.3),
                 desc, font=F_GOTHIC, size=9, color=C_MUTED)
        add_hline(slide, Inches(0.7), y + Inches(0.85), Inches(8.6))
        y += Inches(1.0)

    add_footer(slide, "所要 約30分")

# ============================================================
# スライド 03: ヒアリング振り返り（お悩み4枚）
# ============================================================
def slide_pains():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  振り返り", "03 / 12")
    add_title_block(slide, "前回ヒアリングで出た課題",
                    "ヒアリング内容を4点に整理")

    pains = [
        ("PAIN 01", "リール投稿の続け方が分からない。\n月1本がやっと"),
        ("PAIN 02", "編集に時間がかかりすぎる。\nテロップ位置に悩む"),
        ("PAIN 03", "「うたかたらしさ」を\nブレずに保ちたい"),
        ("PAIN 04", "「劇的変化」を煽る発信に\n違和感がある"),
    ]
    card_w = Inches(4.1); card_h = Inches(1.8)
    positions = [
        (Inches(0.7), Inches(2.5)),
        (Inches(5.2), Inches(2.5)),
        (Inches(0.7), Inches(4.5)),
        (Inches(5.2), Inches(4.5)),
    ]
    for (x, y), (label, text) in zip(positions, pains):
        add_rect(slide, x, y, card_w, card_h, fill=C_CARD, line=C_LINE)
        # 左ボーダー
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            x, y, Inches(0.06), card_h)
        accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = C_ACCENT
        accent_bar.line.fill.background()
        add_text(slide, x + Inches(0.25), y + Inches(0.2), Inches(2), Inches(0.3),
                 label, font=F_MONO, size=9, color=C_ACCENT, bold=True)
        add_text(slide, x + Inches(0.25), y + Inches(0.6), card_w - Inches(0.4), Inches(1.2),
                 text, font=F_MINCHO, size=13, color=C_TEXT, line_spacing=1.5)

    add_footer(slide, "※ ヒアリングメモから抜粋")

# ============================================================
# スライド 04: 3つの軸
# ============================================================
def slide_axes():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  振り返り", "04 / 12")
    add_title_block(slide, "ブランドの軸",
                    "提案の判断基準")

    axes = [
        ("静", "静かに整える", "派手な表現ではなく、\n呼吸が落ち着くような\n言葉と画"),
        ("間", "押し売りしない", "来てください、より\n「そういう時間がある」と\nそっと置いておく"),
        ("続", "無理なく続ける", "毎日張り切る必要はない\n月3〜4本のリズムで\n長く続ける"),
    ]
    card_w = Inches(2.7); card_h = Inches(3.5)
    start_x = Inches(0.85)
    gap = Inches(0.18)
    y = Inches(2.5)
    for i, (icon, title, desc) in enumerate(axes):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, y, card_w, card_h, fill=C_CARD, line=C_LINE)
        # 大きな漢字アイコン
        add_text(slide, x, y + Inches(0.45), card_w, Inches(1.0),
                 icon, font=F_MINCHO, size=44, color=C_ACCENT,
                 align=PP_ALIGN.CENTER)
        # タイトル
        add_text(slide, x, y + Inches(1.6), card_w, Inches(0.5),
                 title, font=F_MINCHO, size=16, color=C_PRIMARY_DK,
                 bold=True, align=PP_ALIGN.CENTER)
        # 説明
        add_text(slide, x + Inches(0.2), y + Inches(2.25), card_w - Inches(0.4), Inches(1.2),
                 desc, font=F_GOTHIC, size=10, color=C_MUTED,
                 align=PP_ALIGN.CENTER, line_spacing=1.6)

    add_footer(slide, "提案はこの軸に沿って組み立て", accent=True)

# ============================================================
# スライド 05: 解決の方向性
# ============================================================
def slide_approach():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  解決の方向性", "05 / 12")
    add_title_block(slide, "3つのアプローチ",
                    "AI任せにせず、人の判断を残す設計")

    steps = [
        ("① 言葉づくり", "ChatGPTを\n道具として使う\nうたかた専用プロンプト"),
        ("② 動画編集",   "テロップ位置を\n悩まないツール\nブラウザで完結"),
        ("③ 撮影設計",   "1シーン1動画で\n短く分けて撮影\n編集が10分に"),
    ]
    card_w = Inches(2.5); card_h = Inches(2.3)
    arrow_w = Inches(0.4)
    start_x = Inches(0.6)
    y = Inches(2.7)
    for i, (title, desc) in enumerate(steps):
        x = start_x + i * (card_w + arrow_w)
        add_rect(slide, x, y, card_w, card_h, fill=C_CARD, line=C_LINE)
        add_text(slide, x, y + Inches(0.4), card_w, Inches(0.5),
                 title, font=F_MINCHO, size=14, color=C_PRIMARY_DK,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), y + Inches(1.05), card_w - Inches(0.3), Inches(1.2),
                 desc, font=F_GOTHIC, size=10, color=C_MUTED,
                 align=PP_ALIGN.CENTER, line_spacing=1.6)
        if i < 2:
            add_text(slide, x + card_w, y + Inches(0.95), arrow_w, Inches(0.5),
                     "▶", font=F_GOTHIC, size=16, color=C_ACCENT,
                     align=PP_ALIGN.CENTER)
    # 補足
    add_text(slide, Inches(0.7), Inches(5.4), Inches(8.6), Inches(0.6),
             "本日は ① と ② を扱う。\n③（撮影設計）は別資料を参照。",
             font=F_MINCHO, size=12, color=C_TEXT, line_spacing=1.7)

    add_footer(slide, "参考資料:  salon-highlight-guide.html")

# ============================================================
# スライド 06: 提案①バナー
# ============================================================
def slide_proposal_1():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  提案 ①", "06 / 12")

    # バナー
    bx, by = Inches(0.7), Inches(1.2)
    bw, bh = Inches(8.6), Inches(2.0)
    add_rect(slide, bx, by, bw, bh, fill=C_DEEP, line=None)
    add_text(slide, bx + Inches(0.5), by + Inches(0.35), bw, Inches(0.35),
             "PROPOSAL 01", font=F_MONO, size=10, color=C_ACCENT, bold=True)
    add_text(slide, bx + Inches(0.5), by + Inches(0.75), bw, Inches(0.7),
             "プロンプト一覧", font=F_MINCHO, size=26, color=C_WHITE)
    add_text(slide, bx + Inches(0.5), by + Inches(1.45), bw, Inches(0.4),
             "ChatGPTで投稿の言葉づくりを高速化",
             font=F_GOTHIC, size=11, color=C_SUBTLE)

    # 本文
    add_text(slide, Inches(0.7), Inches(3.7), Inches(8.6), Inches(2.5),
             "ChatGPTに「リール5本」のような短いトリガー語を打つと、\n"
             "ブランドに沿った投稿案がまとめて出力される仕組み。\n\n"
             "投稿のたびにゼロから考える時間を削減。",
             font=F_MINCHO, size=14, color=C_TEXT, line_spacing=1.9)

    add_footer(slide, "実物: salon-prompts.html  /  15種類のトリガー（次ページ）")

# ============================================================
# スライド 07: 15トリガー
# ============================================================
def slide_triggers():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  提案 ①", "07 / 12")
    add_title_block(slide, "15種類のトリガー",
                    "用途別のトリガー語を15種類用意")

    triggers = [
        ("リール5本", "REEL"), ("フィード3本", "FEED"), ("ストーリーズ", "STORY"),
        ("キャプション", "CAPTION"), ("ハッシュタグ", "HASH"),
        ("月次テーマ", "MONTH"), ("週次テーマ", "WEEK"), ("部位別投稿", "PART"),
        ("悩み言語化", "PAIN"), ("考え方を伝える", "VOICE"),
        ("予約導線", "CTA"), ("LINE配信", "LINE"), ("お客様の声", "VOC"),
        ("季節投稿", "SEASON"), ("NGチェック", "NG"),
    ]
    cols = 5
    cw = Inches(1.65); ch = Inches(0.9)
    gap_x = Inches(0.05); gap_y = Inches(0.15)
    start_x = Inches(0.75)
    start_y = Inches(2.5)
    for i, (jp, en) in enumerate(triggers):
        row = i // cols
        col = i % cols
        x = start_x + col * (cw + gap_x)
        y = start_y + row * (ch + gap_y)
        add_rect(slide, x, y, cw, ch, fill=C_CARD, line=C_LINE)
        add_text(slide, x, y + Inches(0.18), cw, Inches(0.35),
                 jp, font=F_MINCHO, size=11, color=C_PRIMARY_DK,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + Inches(0.52), cw, Inches(0.25),
                 en, font=F_MONO, size=8, color=C_ACCENT,
                 align=PP_ALIGN.CENTER, bold=True)

    add_footer(slide, "例: 月の方針が定まらない時は「月次テーマ」", accent=True)

# ============================================================
# スライド 08: 使い方フロー
# ============================================================
def slide_prompt_flow():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  提案 ①", "08 / 12")
    add_title_block(slide, "使い方は3ステップ",
                    "事前準備は1回のみ")

    steps = [
        ("1", "事前準備",   "ChatGPTに\nプロジェクト用\nプロンプトを貼る\n（最初の1回だけ）"),
        ("2", "毎日の使用", "トリガー語を\n1〜2語打つ\n例:「リール5本」"),
        ("3", "微調整",     "出てきた案から\n使える1つを選んで\n言葉を少し直す"),
    ]
    card_w = Inches(2.5); card_h = Inches(3.0)
    arrow_w = Inches(0.4)
    start_x = Inches(0.6)
    y = Inches(2.7)
    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (card_w + arrow_w)
        add_rect(slide, x, y, card_w, card_h, fill=C_CARD, line=C_LINE)
        # 番号丸
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            x + (card_w - Inches(0.5)) / 2, y + Inches(0.35),
            Inches(0.5), Inches(0.5))
        circle.fill.solid(); circle.fill.fore_color.rgb = C_PRIMARY
        circle.line.fill.background()
        add_text(slide, x, y + Inches(0.4), card_w, Inches(0.4),
                 num, font=F_MONO, size=14, color=C_WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
        # タイトル
        add_text(slide, x, y + Inches(1.05), card_w, Inches(0.4),
                 title, font=F_MINCHO, size=14, color=C_PRIMARY_DK,
                 bold=True, align=PP_ALIGN.CENTER)
        # 説明
        add_text(slide, x + Inches(0.15), y + Inches(1.6), card_w - Inches(0.3), Inches(1.3),
                 desc, font=F_GOTHIC, size=10, color=C_MUTED,
                 align=PP_ALIGN.CENTER, line_spacing=1.6)
        if i < 2:
            add_text(slide, x + card_w, y + Inches(1.25), arrow_w, Inches(0.5),
                     "▶", font=F_GOTHIC, size=16, color=C_ACCENT,
                     align=PP_ALIGN.CENTER)
    # フッターメッセージ
    add_text(slide, 0, Inches(6.15), SW, Inches(0.4),
             "1投稿あたり 3〜5分 で完成",
             font=F_GOTHIC, size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, "本日デモを実施")

# ============================================================
# スライド 09: 提案②バナー
# ============================================================
def slide_proposal_2():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  提案 ②", "09 / 12")

    # バナー
    bx, by = Inches(0.7), Inches(1.2)
    bw, bh = Inches(8.6), Inches(2.0)
    add_rect(slide, bx, by, bw, bh, fill=C_DEEP, line=None)
    add_text(slide, bx + Inches(0.5), by + Inches(0.35), bw, Inches(0.35),
             "PROPOSAL 02", font=F_MONO, size=10, color=C_ACCENT, bold=True)
    add_text(slide, bx + Inches(0.5), by + Inches(0.75), bw, Inches(0.7),
             "動画エディタ", font=F_MINCHO, size=26, color=C_WHITE)
    add_text(slide, bx + Inches(0.5), by + Inches(1.45), bw, Inches(0.4),
             "テロップ位置を悩まないブラウザ完結ツール",
             font=F_GOTHIC, size=11, color=C_SUBTLE)

    # 本文
    add_text(slide, Inches(0.7), Inches(3.7), Inches(8.6), Inches(2.8),
             "動画にテロップを重ねるのではなく、\n"
             "動画の上下に固定の枠を設けて文字を入れる仕組み。\n\n"
             "・テロップが顔に被らない\n"
             "・位置調整が不要\n"
             "・ブランドの統一感を保ったまま量産可能",
             font=F_MINCHO, size=13, color=C_TEXT, line_spacing=1.8)

    add_footer(slide, "アクセス: utakata-subtitle.vercel.app/salon-editor.html")

# ============================================================
# スライド 10: エディタの仕組み
# ============================================================
def slide_editor_visual():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  提案 ②", "10 / 12")
    add_title_block(slide, "エディタの仕組み",
                    "上下の固定枠と中央の動画エリアの3層構造")

    # 左: モック (9:16)
    mock_x = Inches(0.9); mock_y = Inches(2.5)
    mock_w = Inches(2.1); mock_h = Inches(3.8)  # 9:16 比率近似
    # フレーム外枠
    add_rect(slide, mock_x, mock_y, mock_w, mock_h, fill=C_WHITE, line=C_LINE)
    # 上枠
    top_h = mock_h * 0.146
    add_rect(slide, mock_x, mock_y, mock_w, top_h, fill=C_BG, line=None)
    add_text(slide, mock_x, mock_y + Inches(0.12), mock_w, Inches(0.3),
             "呼吸が浅い日に", font=F_MINCHO, size=11, color=C_TEXT,
             bold=True, align=PP_ALIGN.CENTER)
    # 中央 (動画)
    mid_y = mock_y + top_h
    mid_h = mock_h * (1 - 0.292)
    add_rect(slide, mock_x, mid_y, mock_w, mid_h, fill=RGBColor(0x2A, 0x2A, 0x2A), line=None)
    add_text(slide, mock_x, mid_y + mid_h/2 - Inches(0.15), mock_w, Inches(0.3),
             "▶ 動画", font=F_GOTHIC, size=10, color=C_SUBTLE, align=PP_ALIGN.CENTER)
    # 下枠
    bot_y = mid_y + mid_h
    bot_h = mock_h * 0.146
    add_rect(slide, mock_x, bot_y, mock_w, bot_h, fill=C_ACCENT, line=None)
    add_text(slide, mock_x, bot_y + Inches(0.12), mock_w, Inches(0.3),
             "UTAKATA / DM", font=F_GOTHIC, size=10, color=C_WHITE,
             bold=True, align=PP_ALIGN.CENTER)

    # 右: 機能リスト
    feats = [
        ("上の枠（明朝）に主題テロップ",
         "例:「呼吸が浅い日に」「自分に戻る時間」"),
        ("真ん中に動画",
         "撮影素材をそのまま入力"),
        ("下の枠（ゴシック）にCTA",
         "例:「UTAKATA / ご予約はDMから」"),
        ("時間で切り替え可能",
         "0〜15秒は文言A、15〜30秒は文言B"),
    ]
    fx = Inches(3.6); fy = Inches(2.5); fw = Inches(5.7)
    for title, desc in feats:
        add_text(slide, fx, fy, Inches(0.3), Inches(0.4),
                 "○", font=F_MINCHO, size=14, color=C_ACCENT, bold=True)
        add_text(slide, fx + Inches(0.35), fy, fw, Inches(0.4),
                 title, font=F_MINCHO, size=13, color=C_PRIMARY_DK, bold=True)
        add_text(slide, fx + Inches(0.35), fy + Inches(0.4), fw, Inches(0.4),
                 desc, font=F_GOTHIC, size=10, color=C_MUTED)
        fy += Inches(0.95)

    add_footer(slide, "テロップが動画に被らない構造", accent=True)

# ============================================================
# スライド 11: エディタ4ステップ
# ============================================================
def slide_editor_steps():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  提案 ②", "11 / 12")
    add_title_block(slide, "使い方は4ステップ",
                    "本日デモで実演")

    steps = [
        ("STEP 1", "動画を選ぶ", "スマホで撮った\n動画ファイルを1本選択。\n30秒前後が目安。"),
        ("STEP 2", "枠を選ぶ",   "和紙風 / コントラスト /\nクレイの3種から選択。"),
        ("STEP 3", "テロップを入れる", "上枠・下枠それぞれ、\n時間と文字を入力。\nチップで定型文も。"),
        ("STEP 4", "書き出し",   "ボタン1つでMP4が出力。\nそのままインスタに\nアップ可能。"),
    ]
    cw = Inches(2.0); ch = Inches(2.7)
    gap = Inches(0.15)
    start_x = Inches(0.55)
    y = Inches(2.5)
    for i, (tag, title, desc) in enumerate(steps):
        x = start_x + i * (cw + gap)
        add_rect(slide, x, y, cw, ch, fill=C_CARD, line=C_LINE)
        # 上の太線
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.05))
        top_bar.fill.solid(); top_bar.fill.fore_color.rgb = C_PRIMARY
        top_bar.line.fill.background()
        # タグ
        add_text(slide, x + Inches(0.2), y + Inches(0.25), cw, Inches(0.3),
                 tag, font=F_MONO, size=9, color=C_ACCENT, bold=True)
        # タイトル
        add_text(slide, x + Inches(0.2), y + Inches(0.65), cw - Inches(0.3), Inches(0.5),
                 title, font=F_MINCHO, size=14, color=C_PRIMARY_DK, bold=True)
        # 説明
        add_text(slide, x + Inches(0.2), y + Inches(1.3), cw - Inches(0.3), Inches(1.5),
                 desc, font=F_GOTHIC, size=10, color=C_MUTED, line_spacing=1.6)

    add_text(slide, 0, Inches(5.55), SW, Inches(0.4),
             "1本あたり目安 10分",
             font=F_GOTHIC, size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, "※ アプリのインストール不要。スマホのSafari/Chromeで使えます")

# ============================================================
# スライド 12: 3週間ロードマップ
# ============================================================
def slide_roadmap():
    slide = add_blank_slide()
    add_header(slide, "UTAKATA  —  次のアクション", "12 / 12")
    add_title_block(slide, "次回までの3週間ロードマップ",
                    "段階的に導入")

    weeks = [
        ("WEEK 01", "プロンプトを試す",
         "・ChatGPTにプロジェクト用プロンプトを設定\n・「リール5本」を1回実行\n・出力から1案を選定"),
        ("WEEK 02", "エディタで1本作る",
         "・撮影素材を1本用意\n・エディタで枠付き動画化\n・書き出してインスタに投稿"),
        ("WEEK 03", "運用に乗せる",
         "・もう1〜2本作成\n・使いにくい部分をメモ\n・次回ヒアリングで共有"),
    ]
    card_w = Inches(2.7); card_h = Inches(3.4)
    gap = Inches(0.2)
    start_x = Inches(0.85)
    y = Inches(2.5)
    for i, (tag, title, desc) in enumerate(weeks):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, y, card_w, card_h, fill=C_CARD, line=C_LINE)
        add_text(slide, x + Inches(0.25), y + Inches(0.3), card_w, Inches(0.3),
                 tag, font=F_MONO, size=10, color=C_ACCENT, bold=True)
        add_text(slide, x + Inches(0.25), y + Inches(0.7), card_w - Inches(0.5), Inches(0.5),
                 title, font=F_MINCHO, size=15, color=C_PRIMARY_DK, bold=True)
        add_text(slide, x + Inches(0.25), y + Inches(1.4), card_w - Inches(0.5), Inches(1.8),
                 desc, font=F_GOTHIC, size=10, color=C_TEXT, line_spacing=1.7)

    add_footer(slide, "次回ヒアリングで運用状況を確認", accent=True)

# ============================================================
# スライド 13: 終わり
# ============================================================
def slide_closing():
    slide = add_blank_slide(C_DEEP)
    add_text(slide, 0, Inches(3.0), SW, Inches(1.0),
             "以上", font=F_MINCHO, size=36, color=C_WHITE,
             align=PP_ALIGN.CENTER)
    # 装飾ライン
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(4.7), Inches(4.2), Inches(0.6), Emu(9525))
    deco.fill.solid(); deco.fill.fore_color.rgb = C_ACCENT
    deco.line.fill.background()
    add_text(slide, 0, Inches(4.7), SW, Inches(0.6),
             "ご質問・ご意見をお願いします",
             font=F_GOTHIC, size=12, color=C_SUBTLE,
             align=PP_ALIGN.CENTER)

# ===== 実行 =====
slide_cover()
slide_agenda()
slide_pains()
slide_axes()
slide_approach()
slide_proposal_1()
slide_triggers()
slide_prompt_flow()
slide_proposal_2()
slide_editor_visual()
slide_editor_steps()
slide_roadmap()
slide_closing()

prs.save('/home/user/utakata-subtitle/salon-slides.pptx')
print(f"Saved: {len(prs.slides)} slides")
